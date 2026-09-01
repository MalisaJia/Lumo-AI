"""会话导出核心逻辑：生成 PDF（fpdf2）与 PPTX（python-pptx）。

只读取 Conversation/Message ORM 对象与 uploads 目录，不修改任何数据。
"""

import json
import logging
import os
from datetime import datetime, timezone
from io import BytesIO
from pathlib import Path

from fpdf import FPDF
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from app.models import Conversation, Message

logger = logging.getLogger(__name__)

# 单条消息最大字符数（超出截断并标注）
MAX_MESSAGE_CHARS = 10000
# PPTX 幻灯片空间有限，单条消息展示上限更严格
PPTX_MESSAGE_CHARS = 1200

# 主题色（violet-600 / blue-500）
VIOLET = (124, 58, 237)  # #7C3AED
BLUE = (59, 130, 246)  # #3B82F6
GRAY = (107, 114, 128)  # 辅助灰
DARK = (31, 41, 55)  # 正文深灰

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".gif"}

# Windows 系统 CJK 字体候选（(常规, 加粗)，按优先级排列）
_FONT_CANDIDATES: list[tuple[str, str | None]] = [
    ("msyh.ttc", "msyhbd.ttc"),  # 微软雅黑
    ("msyh.ttf", "msyhbd.ttf"),
    ("Deng.ttf", "Dengb.ttf"),  # 等线
    ("simhei.ttf", None),  # 黑体
    ("simsun.ttc", None),  # 宋体
]


# ---------------------------------------------------------------- 通用辅助


def _parse_json_list(text: str | None) -> list[dict]:
    """解析 sources/attachments JSON 字符串；脏数据一律返回空列表。"""
    if not text:
        return []
    try:
        data = json.loads(text)
    except (ValueError, TypeError):
        return []
    return [item for item in data if isinstance(item, dict)] if isinstance(data, list) else []


def _truncate(content: str, limit: int = MAX_MESSAGE_CHARS) -> str:
    if len(content) > limit:
        return content[:limit] + "\n[内容已截断]"
    return content


def _role_label(role: str) -> str:
    return {"user": "用户", "assistant": "助手"}.get(role, role)


def _role_color(role: str) -> tuple[int, int, int]:
    return BLUE if role == "user" else VIOLET


def _fmt_time(dt: datetime | None) -> str:
    """UTC 存储的 naive 时间转本地时区展示。"""
    if not dt:
        return ""
    if dt.tzinfo is None:
        dt = dt.replace(tzinfo=timezone.utc)
    return dt.astimezone().strftime("%Y-%m-%d %H:%M")


def _image_attachments(message: Message, upload_dir: Path) -> list[tuple[str, Path | None]]:
    """从 attachments JSON 提取图片 (文件名, 本地路径或 None-缺失)。"""
    result: list[tuple[str, Path | None]] = []
    for item in _parse_json_list(message.attachments):
        url = str(item.get("url") or "")
        name = url.rsplit("/", 1)[-1]
        if not name or Path(name).suffix.lower() not in IMAGE_EXTS:
            continue
        path = upload_dir / name
        result.append((name, path if path.is_file() else None))
    return result


def _message_sources(message: Message) -> list[dict]:
    return [s for s in _parse_json_list(message.sources) if s.get("url") or s.get("title")]


# ---------------------------------------------------------------- PDF 生成


def _find_cjk_fonts() -> tuple[str | None, str | None]:
    """在系统字体目录查找 CJK 字体，返回 (常规路径, 加粗路径或 None)。"""
    fonts_dir = Path(os.environ.get("WINDIR", r"C:\Windows")) / "Fonts"
    for regular, bold in _FONT_CANDIDATES:
        reg_path = fonts_dir / regular
        if reg_path.is_file():
            bold_path = fonts_dir / bold if bold else None
            return str(reg_path), (
                str(bold_path) if bold_path and bold_path.is_file() else None
            )
    return None, None


def _register_cjk_font(pdf: FPDF) -> tuple[str, bool]:
    """注册 CJK 字体，返回 (字体族名, 是否支持 Unicode)。

    找不到系统字体或注册失败时回退 Helvetica（CJK 字符会被替换为 ?）。
    """
    regular, bold = _find_cjk_fonts()
    if regular:
        try:
            pdf.add_font("cjk", "", regular)
            # 无独立粗体文件时用常规字重充当粗体样式，保证 style="B" 可用
            pdf.add_font("cjk", "B", bold or regular)
            return "cjk", True
        except Exception:
            logger.warning("注册 CJK 字体失败：%s", regular, exc_info=True)
    return "helvetica", False


def _safe_text(text: str, unicode_ok: bool) -> str:
    """Helvetica 回退时清洗 latin-1 之外的字符，避免编码异常。"""
    if unicode_ok:
        return text
    return text.encode("latin-1", "replace").decode("latin-1")


def _pdf_embed_image(pdf: FPDF, path: Path) -> bool:
    """按页宽 80% 比例居中嵌入图片；失败返回 False。"""
    try:
        with Image.open(path) as img:
            img.load()
            width = pdf.epw * 0.8
            height = width * img.height / img.width
            # 超过整页可用高度时按高度反向约束
            max_h = pdf.h - pdf.t_margin - pdf.b_margin
            if height > max_h:
                width = width * max_h / height
                height = max_h
            if pdf.get_y() + height > pdf.page_break_trigger:
                pdf.add_page()
            x = pdf.l_margin + (pdf.epw - width) / 2
            pdf.image(img, x=x, y=pdf.get_y(), w=width, h=height)
        pdf.set_y(pdf.get_y() + height + 2)
        return True
    except Exception:
        logger.warning("PDF 嵌入图片失败：%s", path, exc_info=True)
        return False


def generate_pdf(
    conversation: Conversation, messages: list[Message], upload_dir: Path
) -> BytesIO:
    """将会话导出为 PDF，返回文件内容的 BytesIO。"""
    pdf = FPDF(orientation="P", unit="mm", format="A4")
    pdf.set_margins(15, 15, 15)
    pdf.set_auto_page_break(auto=True, margin=15)
    font, unicode_ok = _register_cjk_font(pdf)
    pdf.add_page()

    # 首页头部：标题 + 导出时间 + 消息总数
    pdf.set_font(font, style="B", size=18)
    pdf.set_text_color(*DARK)
    pdf.multi_cell(0, 10, _safe_text(conversation.title, unicode_ok), align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(2)
    pdf.set_font(font, size=10)
    pdf.set_text_color(*GRAY)
    meta = f"导出时间：{_fmt_time(datetime.now(timezone.utc))}    消息总数：{len(messages)}"
    pdf.cell(0, 6, _safe_text(meta, unicode_ok), align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(6)

    for message in messages:
        # 角色标签（加粗着色）+ 时间（灰色 8pt）
        pdf.set_font(font, style="B", size=11)
        pdf.set_text_color(*_role_color(message.role))
        label = _safe_text(_role_label(message.role), unicode_ok)
        pdf.cell(pdf.get_string_width(label) + 2, 6, label)
        pdf.set_font(font, size=8)
        pdf.set_text_color(*GRAY)
        pdf.cell(0, 6, _fmt_time(message.created_at), new_x="LMARGIN", new_y="NEXT")
        pdf.ln(1)

        # 正文内容（10pt 黑色，超长截断）
        pdf.set_font(font, size=10)
        pdf.set_text_color(0, 0, 0)
        content = _safe_text(_truncate(message.content or ""), unicode_ok)
        if content.strip():
            pdf.multi_cell(0, 5.5, content, new_x="LMARGIN", new_y="NEXT")

        # 图片附件：缺失或解码失败写占位文本
        for name, path in _image_attachments(message, upload_dir):
            pdf.ln(1)
            if path is None or not _pdf_embed_image(pdf, path):
                pdf.set_font(font, size=9)
                pdf.set_text_color(*GRAY)
                pdf.cell(
                    0,
                    5,
                    _safe_text(f"[图片缺失] {name}", unicode_ok),
                    new_x="LMARGIN",
                    new_y="NEXT",
                )

        # 参考来源（小字灰色）
        sources = _message_sources(message)
        if sources:
            pdf.ln(1)
            pdf.set_font(font, size=8)
            pdf.set_text_color(*GRAY)
            for idx, source in enumerate(sources, 1):
                line = f"参考：[{idx}] {source.get('title') or ''} ({source.get('url') or ''})"
                # multi_cell 默认 new_x=RIGHT，连续调用需显式回到左边距
                pdf.multi_cell(0, 4.5, _safe_text(line, unicode_ok), new_x="LMARGIN", new_y="NEXT")

        # 水平分隔线
        pdf.ln(2)
        y = pdf.get_y()
        if y < pdf.page_break_trigger:
            pdf.set_draw_color(229, 231, 235)
            pdf.line(pdf.l_margin, y, pdf.w - pdf.r_margin, y)
        pdf.ln(4)

    return BytesIO(bytes(pdf.output()))


# ---------------------------------------------------------------- PPTX 生成

_SLIDE_W = Inches(13.333)  # 16:9 宽屏
_SLIDE_H = Inches(7.5)


def _pptx_message_block(slide, message: Message, top: float) -> None:
    """在幻灯片指定纵向位置渲染一条消息（角色 + 时间 + 内容）。"""
    box = slide.shapes.add_textbox(
        Inches(0.6), Inches(top), _SLIDE_W - Inches(1.2), Inches(3.1)
    )
    tf = box.text_frame
    tf.word_wrap = True

    head = tf.paragraphs[0]
    run = head.add_run()
    run.text = _role_label(message.role)
    run.font.bold = True
    run.font.size = Pt(16)
    run.font.color.rgb = RGBColor(*_role_color(message.role))
    time_run = head.add_run()
    time_run.text = f"   {_fmt_time(message.created_at)}"
    time_run.font.size = Pt(10)
    time_run.font.color.rgb = RGBColor(*GRAY)

    body = tf.add_paragraph()
    content = _truncate(message.content or "", PPTX_MESSAGE_CHARS)
    body.text = content if content.strip() else "（图片消息）"
    body.font.size = Pt(13)
    body.font.color.rgb = RGBColor(*DARK)


def _pptx_image_slide(prs, name: str, path: Path | None) -> None:
    """图片独立幻灯片：经 PIL 统一转 PNG（兼容 webp 等格式）并居中。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if path is not None:
        try:
            with Image.open(path) as img:
                img.load()
                buf = BytesIO()
                img.convert("RGBA").save(buf, format="PNG")
                buf.seek(0)
                max_w, max_h = _SLIDE_W - Inches(1.6), _SLIDE_H - Inches(1.2)
                scale = min(max_w / img.width, max_h / img.height)
                width, height = int(img.width * scale), int(img.height * scale)
            slide.shapes.add_picture(
                buf,
                left=int((_SLIDE_W - width) / 2),
                top=int((_SLIDE_H - height) / 2),
                width=width,
                height=height,
            )
            return
        except Exception:
            logger.warning("PPTX 嵌入图片失败：%s", path, exc_info=True)
    box = slide.shapes.add_textbox(
        Inches(0.6), Inches(3.3), _SLIDE_W - Inches(1.2), Inches(1)
    )
    para = box.text_frame.paragraphs[0]
    para.text = f"[图片缺失] {name}"
    para.font.size = Pt(14)
    para.font.color.rgb = RGBColor(*GRAY)


def generate_pptx(
    conversation: Conversation, messages: list[Message], upload_dir: Path
) -> BytesIO:
    """将会话导出为 16:9 PPTX，返回文件内容的 BytesIO。"""
    prs = Presentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H

    # 标题幻灯片：会话标题 + 副标题（日期 + 模型名）
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = conversation.title
    for para in title_slide.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(*VIOLET)
    subtitle_parts = [_fmt_time(datetime.now(timezone.utc))]
    if conversation.model_name:
        subtitle_parts.append(conversation.model_name)
    subtitle = title_slide.placeholders[1]
    subtitle.text = " · ".join(subtitle_parts)
    for para in subtitle.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(*BLUE)

    # 内容幻灯片：纯文本消息每张最多 2 条；带图片的消息拆出独立图片页
    text_queue: list[Message] = []

    def flush_text_queue() -> None:
        while text_queue:
            batch = text_queue[:2]
            del text_queue[:2]
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            for i, msg in enumerate(batch):
                _pptx_message_block(slide, msg, top=0.4 + i * 3.6)

    for message in messages:
        images = _image_attachments(message, upload_dir)
        text_queue.append(message)
        if images:
            flush_text_queue()
            for name, path in images:
                _pptx_image_slide(prs, name, path)
    flush_text_queue()

    # 引用汇总幻灯片：全部 sources 按 url 去重，每页最多 12 条
    seen: set[str] = set()
    all_sources: list[dict] = []
    for message in messages:
        for source in _message_sources(message):
            key = source.get("url") or source.get("title") or ""
            if key and key not in seen:
                seen.add(key)
                all_sources.append(source)
    for start in range(0, len(all_sources), 12):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        head_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(0.4), _SLIDE_W - Inches(1.2), Inches(0.8)
        )
        head = head_box.text_frame.paragraphs[0]
        head.text = "参考来源"
        head.font.bold = True
        head.font.size = Pt(22)
        head.font.color.rgb = RGBColor(*VIOLET)
        body_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.3), _SLIDE_W - Inches(1.2), _SLIDE_H - Inches(1.7)
        )
        tf = body_box.text_frame
        tf.word_wrap = True
        for offset, source in enumerate(all_sources[start : start + 12]):
            para = tf.paragraphs[0] if offset == 0 else tf.add_paragraph()
            para.text = (
                f"[{start + offset + 1}] {source.get('title') or ''}"
                f" — {source.get('url') or ''}"
            )
            para.font.size = Pt(12)
            para.font.color.rgb = RGBColor(*DARK)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
